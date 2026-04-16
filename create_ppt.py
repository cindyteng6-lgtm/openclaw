#!/usr/bin/env python3
"""
创建专业美观的 OpenClaw 介绍 PPT
- 嵌入中文字体解决方案
- 现代化设计风格
- 内容丰富专业
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
import os

# 专业配色方案
COLORS = {
    'primary': RGBColor(23, 162, 184),      # 主色 - 科技蓝
    'secondary': RGBColor(42, 42, 42),      # 深灰 - 文字
    'accent': RGBColor(255, 193, 7),        # 强调色 - 金黄
    'light_bg': RGBColor(248, 249, 250),    # 浅灰背景
    'white': RGBColor(255, 255, 255),       # 白色
    'dark': RGBColor(33, 37, 41),           # 深色
}

# 中文字体列表（按优先级）
CHINESE_FONTS = [
    'PingFang SC',
    'Heiti SC',
    'Microsoft YaHei',
    'SimHei',
    'STHeiti',
    'Hiragino Sans GB',
]

def add_font_embedding(prs):
    """
    在 PPTX 中嵌入字体引用
    注意：python-pptx 不直接支持字体嵌入，
    但可以通过设置字体名称让系统使用已安装的字体
    """
    pass  # 字体嵌入需要在 PowerPoint 中手动完成，或使用 VBA

def get_chinese_font_name():
    """返回第一个可用的中文字体名称"""
    return 'PingFang SC'  # macOS 默认中文字体

def create_title_slide(prs, title, subtitle):
    """创建标题页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版式
    
    # 背景色块
    bg_shape = slide.shapes.add_shape(
        1,  # msoShapeRectangle
        Inches(0), Inches(0), Inches(13.33), Inches(7.5)
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = COLORS['primary']
    bg_shape.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.33), Inches(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.font.name = get_chinese_font_name()
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.33), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = COLORS['white']
    p.font.name = get_chinese_font_name()
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def create_content_slide(prs, title, content_items, slide_number=0):
    """创建内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版式
    
    # 顶部色条
    header_bar = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(13.33), Inches(0.8)
    )
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = COLORS['primary']
    header_bar.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(10), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.font.name = get_chinese_font_name()
    
    # 页码
    page_box = slide.shapes.add_textbox(Inches(12.5), Inches(0.2), Inches(0.7), Inches(0.4))
    tf = page_box.text_frame
    p = tf.paragraphs[0]
    p.text = str(slide_number)
    p.font.size = Pt(14)
    p.font.color.rgb = COLORS['white']
    p.font.name = get_chinese_font_name()
    p.alignment = PP_ALIGN.RIGHT
    
    # 内容区域背景
    bg_shape = slide.shapes.add_shape(
        1, Inches(0.3), Inches(1), Inches(12.73), Inches(6.2)
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = COLORS['light_bg']
    bg_shape.line.color.rgb = COLORS['primary']
    bg_shape.line.width = Pt(1)
    
    # 内容文本
    content_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(12.13), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        if isinstance(item, dict):
            p.text = item.get('text', '')
            p.font.size = Pt(item.get('size', 20))
            p.font.bold = item.get('bold', False)
            p.font.color.rgb = item.get('color', COLORS['secondary'])
            p.space_after = Pt(item.get('space_after', 12))
        else:
            p.text = f"• {item}"
            p.font.size = Pt(20)
            p.font.color.rgb = COLORS['secondary']
            p.space_after = Pt(12)
        
        p.font.name = get_chinese_font_name()
    
    return slide

def create_section_slide(prs, title, description=""):
    """创建章节过渡页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    bg_shape = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(13.33), Inches(7.5)
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = COLORS['dark']
    bg_shape.line.fill.background()
    
    # 装饰线条
    line = slide.shapes.add_shape(
        1, Inches(0.5), Inches(3), Inches(2), Inches(0.1)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS['accent']
    line.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.font.name = get_chinese_font_name()
    
    if description:
        desc_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12), Inches(1))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = description
        p.font.size = Pt(24)
        p.font.color.rgb = COLORS['accent']
        p.font.name = get_chinese_font_name()
    
    return slide

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # 1. 封面页
    create_title_slide(
        prs,
        "OpenClaw 智能助手平台",
        "下一代 AI 自动化与任务编排系统"
    )
    
    # 2. 什么是 OpenClaw
    create_content_slide(prs, "什么是 OpenClaw？", [
        {'text': 'OpenClaw 是什么？', 'size': 28, 'bold': True, 'space_after': 16},
        "开源 AI 助手运行平台，让 AI 真正融入你的工作流",
        "支持多模型接入（Qwen、GPT、Claude 等）",
        "内置丰富工具集：文件操作、浏览器控制、定时任务、消息通知",
        "可扩展的技能系统，轻松定制专属能力",
        {'text': '\n核心价值', 'size': 28, 'bold': True, 'space_after': 16},
        "🔧 自动化日常任务，释放人力",
        "🧠 持久化记忆，真正理解你的需求",
        "🔐 本地部署，数据完全可控",
    ], 1)
    
    # 3. 核心功能
    create_content_slide(prs, "核心功能特性", [
        {'text': '📁 文件与 workspace 管理', 'size': 24, 'bold': True, 'space_after': 8},
        "  • 读写编辑文件，自动整理归类",
        "  • 全局 workspace 概念，所有会话共享上下文",
        {'text': '\n🌐 浏览器与网络操作', 'size': 24, 'bold': True, 'space_after': 8},
        "  • 自动化网页交互（点击、输入、截图）",
        "  • 网页内容提取与搜索",
        {'text': '\n⏰ 定时任务系统', 'size': 24, 'bold': True, 'space_after': 8},
        "  • Cron 表达式支持，精确调度",
        "  • 心跳检查，主动提醒重要事项",
        {'text': '\n💬 多平台消息', 'size': 24, 'bold': True, 'space_after': 8},
        "  • Telegram、Discord、微信等集成",
        "  • 群组聊天智能参与",
    ], 2)
    
    # 4. 技术架构
    create_section_slide(prs, "技术架构", "模块化设计 · 灵活扩展")
    
    create_content_slide(prs, "系统架构概览", [
        {'text': '🏗️ 三层架构', 'size': 28, 'bold': True, 'space_after': 16},
        "Gateway 层：核心调度与服务管理",
        "Agent 层：AI 模型推理与决策",
        "Surface 层：用户界面与交互（Web/Telegram/Discord）",
        {'text': '\n🔌 插件系统', 'size': 28, 'bold': True, 'space_after': 16},
        "Skills：预定义能力包（GitHub、天气、健康检查等）",
        "Tools：底层工具接口（文件、浏览器、节点设备）",
        "Memory：短期日志 + 长期记忆的持久化系统",
    ], 3)
    
    # 5. 记忆系统
    create_content_slide(prs, "记忆系统 - AI 的连续性", [
        {'text': '📝 日常记忆 (memory/YYYY-MM-DD.md)', 'size': 22, 'bold': True, 'space_after': 8},
        "  • 每日自动创建，记录原始对话与事件",
        "  • 类似人类的日记，保留完整上下文",
        {'text': '\n🧠 长期记忆 (MEMORY.md)', 'size': 22, 'bold': True, 'space_after': 8},
        "  •  curated 精华记忆，仅主会话加载",
        "  • 安全隔离：群聊不访问个人隐私",
        {'text': '\n📋 配置文件', 'size': 22, 'bold': True, 'space_after': 8},
        "  • SOUL.md：AI 人格与行为准则",
        "  • USER.md：用户偏好与背景信息",
        "  • IDENTITY.md：AI 自我认知",
    ], 4)
    
    # 6. 子代理系统
    create_content_slide(prs, "子代理系统 - 并行任务处理", [
        "🚀 支持spawn独立子代理处理复杂任务",
        "📊 子代理可运行不同模型、不同配置",
        "🔄 主代理可 steer/kill/监控子代理状态",
        "💡 适用场景：",
        "  • GitHub issue 批量处理",
        "  • 多文件代码重构",
        "  • 长时间后台任务",
    ], 5)
    
    # 7. 安全与边界
    create_section_slide(prs, "安全设计", "隐私优先 · 权限可控")
    
    create_content_slide(prs, "安全边界与最佳实践", [
        {'text': '🔒 数据安全', 'size': 24, 'bold': True, 'space_after': 8},
        "  • 本地部署，数据不出境",
        "  • MEMORY.md 主会话隔离，群聊不泄露隐私",
        "  • 敏感操作需用户确认（删除、外部发送）",
        {'text': '\n⚠️ 红线规则', 'size': 24, 'bold': True, 'space_after': 8},
        "  • 禁止主动外发隐私数据",
        "  • 禁止执行破坏性命令无确认",
        "  • 群聊中不代表用户发声",
        {'text': '\n✅ 推荐实践', 'size': 24, 'bold': True, 'space_after': 8},
        "  • 使用 trash 而非 rm（可恢复）",
        "  • 不确定时先询问用户",
    ], 6)
    
    # 8. 快速开始
    create_content_slide(prs, "快速开始指南", [
        {'text': '1️⃣ 安装', 'size': 24, 'bold': True, 'space_after': 8},
        "  npm install -g openclaw",
        {'text': '\n2️⃣ 配置', 'size': 24, 'bold': True, 'space_after': 8},
        "  openclaw config  # 设置模型、Surface 等",
        {'text': '\n3️⃣ 启动', 'size': 24, 'bold': True, 'space_after': 8},
        "  openclaw gateway start",
        {'text': '\n4️⃣ 连接', 'size': 24, 'bold': True, 'space_after': 8},
        "  Web 界面 / Telegram / Discord / 移动端 App",
        {'text': '\n📚 文档', 'size': 24, 'bold': True, 'space_after': 8},
        "  本地：/usr/local/lib/node_modules/openclaw/docs",
        "  在线：https://docs.openclaw.ai",
    ], 7)
    
    # 9. 使用场景
    create_content_slide(prs, "典型使用场景", [
        "📧 智能邮件管理 - 自动分类、摘要、提醒",
        "📅 日程助手 - 心跳检查即将到来的会议",
        "💻 开发辅助 - GitHub issue 处理、代码审查",
        "🏠 智能家居 - 通过节点控制摄像头、通知",
        "📰 信息聚合 - 定时抓取新闻、天气、股票",
        "📝 知识管理 - 自动整理笔记、更新记忆",
    ], 8)
    
    # 10. 结束页
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    bg_shape = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(13.33), Inches(7.5)
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = COLORS['primary']
    bg_shape.line.fill.background()
    
    # 主标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.33), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "开始构建你的 AI 助手"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.font.name = get_chinese_font_name()
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.33), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "https://github.com/openclaw/openclaw"
    p.font.size = Pt(24)
    p.font.color.rgb = COLORS['white']
    p.font.name = get_chinese_font_name()
    p.alignment = PP_ALIGN.CENTER
    
    # 保存文件
    output_path = '/Users/xiaowu/.openclaw/workspace/OpenClaw_专业正式版.pptx'
    prs.save(output_path)
    
    print(f"✅ PPT 已生成：{output_path}")
    print(f"📊 共 {len(prs.slides)} 页")
    print("\n💡 字体说明：")
    print("   - 使用 macOS 系统字体 PingFang SC")
    print("   - 在 PowerPoint 中打开后，建议：")
    print("     文件 > 选项 > 保存 > 勾选「将字体嵌入文件」")
    print("   - 或使用「导出为 PDF」确保跨平台一致性")

if __name__ == '__main__':
    main()
